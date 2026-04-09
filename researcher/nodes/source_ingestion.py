"""
Input:
- `workspace/input/sources_url.json`
- `workspace/input/source_annotations.json` (optional)
- local files or directories under `workspace/input/`

Output:
- per-source bundles under `workspace/knowledge/sources/<item_id>/`
- `workspace/knowledge/metadata.json`
- `workspace/knowledge/knowledge.md`
"""

from __future__ import annotations

import csv
import hashlib
import json
import os
import re
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.parse import urlparse

import requests
from autogen import ConversableAgent

from researcher.exceptions import WorkflowError
from researcher.prompts.templates import (
    SOURCE_BLOG_PROMPT,
    SOURCE_KNOWLEDGE_SYNTHESIS_PROMPT,
)
from researcher.state import ResearchState
from researcher.utils import (
    assemble_markdown_blog_blocks,
    extract_pdf_markdown_with_images,
    get_llm_config,
    get_relative_path,
    load_artifact_from_file,
    load_markdown,
    resolve_workspace_path,
    rewrite_markdown_paths,
    run_jobs_in_parallel,
    log_stage,
    save_json,
    save_markdown,
)


def source_ingestion_node(state: ResearchState) -> Dict[str, Any]:
    workspace_dir = state["workspace_dir"]
    log_stage(workspace_dir, "source_ingestion", "Starting source ingestion")

    try:
        config = state["config"]["researcher"]["source_ingestion"]
        max_items = config["max_items"]
        max_files_per_item = config["max_files_per_item"]
        max_depth = config["max_depth"]
        max_bytes_per_file = config["max_bytes_per_file"]
        max_total_bytes_per_item = config["max_total_bytes_per_item"]
        max_preview_rows = config["max_preview_rows"]
        timeout_seconds = config["timeout_seconds"]

        # Source ingestion accepts URLs from `input/sources_url.json` plus any
        # user-provided local files or directories placed under `input/`.
        input_dir = workspace_dir / "input"
        if not input_dir.exists() or not input_dir.is_dir():
            raise WorkflowError(f"Input directory not found: {input_dir}")

        sources: List[str] = []

        # Parse explicit URL list from sources_url.json
        sources_url_path = input_dir / "sources_url.json"
        if sources_url_path.exists() and sources_url_path.is_file():
            payload = json.loads(sources_url_path.read_text(encoding="utf-8"))

            url_items = payload.get("sources") or payload.get("urls") or []
            if not isinstance(url_items, list):
                raise WorkflowError("sources_url.json must contain a list under 'sources' or 'urls'")
            for item in url_items:
                if isinstance(item, str) and item.strip():
                    sources.append(item.strip())

        # Parse sources from workspace_dir / "input"
        for entry in sorted(input_dir.iterdir()):
            if entry.name.startswith(".") or entry.name in SOURCE_BLOG_CONTROL_FILES:
                continue
            sources.append(str(entry.resolve()))

        if not sources:
            raise WorkflowError(f"No sources found in input directory: {input_dir}")

        if max_items is not None and len(sources) > max_items:
            sources = sources[:max_items]
            log_stage(
                workspace_dir,
                "source_ingestion",
                f"Source list truncated to max_items={max_items}",
            )

        source_annotations = _load_source_annotations(input_dir=input_dir, workspace_dir=workspace_dir)

        knowledge_dir = workspace_dir / "knowledge"
        bundles_root = knowledge_dir / "sources"
        knowledge_dir.mkdir(parents=True, exist_ok=True)
        bundles_root.mkdir(parents=True, exist_ok=True)

        llm_config = get_llm_config()
        task_context = load_artifact_from_file(workspace_dir, "input") or ""

        # resolve and prepare each source
        metadata_items: List[Dict[str, Any]] = []
        for idx, source in enumerate(sources, start=1):
            item_id = _safe_item_id(idx, source)
            bundle_dir = bundles_root / item_id
            item_meta: Dict[str, Any] = {
                "item_id": item_id,
                "source_input": source,
                "status": "failed",
                "source_type": None,
                "resolved_from": None,
                "local_path": None,
                "bundle_dir": get_relative_path(bundle_dir, workspace_dir),
                "prepared_md_path": None,
                "blog_path": get_relative_path(bundle_dir / "blog.md", workspace_dir),
                "parse_status": "pending",
                "blog_status": "pending",
                "blog_error": None,
                "user_note": "",
                "errors": [],
            }

            annotation = _find_source_annotation(
                source=source,
                annotations=source_annotations,
                workspace_dir=workspace_dir,
                input_dir=input_dir,
            )
            if annotation:
                item_meta["user_note"] = annotation.get("note", "")

            try:
                resolved = _resolve_source_input(
                    source=source,
                    workspace_dir=workspace_dir,
                    bundle_dir=bundle_dir,
                    timeout_seconds=timeout_seconds,
                    max_total_bytes=max_total_bytes_per_item,
                )
                item_meta.update(resolved)
                prepared = _prepare_source_bundle(
                    item_meta=item_meta,
                    workspace_dir=workspace_dir,
                    max_files=max_files_per_item,
                    max_depth=max_depth,
                    max_bytes_per_file=max_bytes_per_file,
                    max_total_bytes_per_item=max_total_bytes_per_item,
                    max_preview_rows=max_preview_rows,
                )
                item_meta.update(prepared)
                item_meta["status"] = "prepared"
                log_stage(workspace_dir, "source_ingestion", f"Prepared source {idx}/{len(sources)}: {source}")
            except Exception as item_error:
                item_meta["errors"].append(str(item_error))
                item_meta["parse_status"] = "failed"
                item_meta["blog_status"] = "failed"
                log_stage(workspace_dir, "source_ingestion", f"Source failed ({source}): {item_error}")

            metadata_items.append(item_meta)

        prepared_items = [item for item in metadata_items if item.get("status") == "prepared"]
        if not prepared_items:
            raise WorkflowError("All source items failed during source ingestion")

        # build blogs in parallel
        log_stage(workspace_dir, "source_ingestion", f"Generating blogs for {len(prepared_items)} prepared sources")
        blog_jobs = [
            {
                "key": idx,
                "prompt": _build_source_blog_prompt(entry, workspace_dir),
                "agent_name": "SourceBlogger",
            }
            for idx, entry in enumerate(prepared_items)
        ]
        indexed_blogs = run_jobs_in_parallel(
            jobs=blog_jobs,
            worker=lambda job: _generate_markdown_with_agent(
                prompt=job["prompt"],
                llm_config=llm_config,
                agent_name=job.get("agent_name", "SourceBlogger"),
            ),
            max_workers=len(blog_jobs),
            progress_desc="Building source blogs",
        )

        for idx, item in enumerate(prepared_items):
            item["status"] = "completed"
            item["blog_status"] = "completed"
            item["summary_path"] = "knowledge/knowledge.md"

        blog_blocks = assemble_markdown_blog_blocks(
            entries=prepared_items,
            workspace_dir=workspace_dir,
            indexed_blogs=indexed_blogs,
            metadata_title="Source Metadata",
            metadata_fields=_source_metadata_field_specs(),
            block_label="Source",
            hint_label="Source Hint",
            hint_builder=_build_source_hint,
            write_blog=True,
        )

        # summarizer
        knowledge_prompt = SOURCE_KNOWLEDGE_SYNTHESIS_PROMPT.format(
            task=task_context,
            blogs_text="\n\n---\n\n".join(blog_blocks),
        )
        knowledge_body = _generate_markdown_with_agent(
            prompt=knowledge_prompt,
            llm_config=llm_config,
            agent_name="SourceKnowledgeSummarizer",
        )
        knowledge_content = knowledge_body.strip()
        metadata_appendix = _build_source_metadata_appendix(prepared_items)
        if metadata_appendix:
            knowledge_content = f"{knowledge_content}\n\n---\n\n{metadata_appendix}"

        success_count = sum(1 for item in metadata_items if item.get("status") == "completed")
        failed_count = len(metadata_items) - success_count
        metadata = {
            "generated_at": datetime.now().isoformat(),
            "node": "source_ingestion",
            "config": config,
            "stats": {
                "total": len(metadata_items),
                "successful": success_count,
                "failed": failed_count,
            },
            "items": metadata_items,
        }

        metadata_path = knowledge_dir / "metadata.json"
        knowledge_path = knowledge_dir / "knowledge.md"
        save_json(metadata, metadata_path)
        save_markdown(knowledge_content, knowledge_path)
        log_stage(workspace_dir, "source_ingestion", f"Wrote knowledge synthesis: {knowledge_path}")

        log_stage(workspace_dir, "source_ingestion", f"Completed. success={success_count}, failed={failed_count}")

        update_state = {
            "stage": "source_ingestion",
            "knowledge": {
                "summary_path": str(knowledge_path),
                "metadata_path": str(metadata_path),
                "stats": metadata.get("stats", {}),
            },
        }
        # router
        if state["config"]["researcher"]["workflow"] == "default":
            update_state["next_node"] = "task_parsing"
        return update_state

    except Exception as e:
        log_stage(workspace_dir, "source_ingestion", f"Error: {str(e)}")
        raise WorkflowError(f"Source ingestion failed: {str(e)}")


TEXT_SUFFIXES = {
    ".md",
    ".txt",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".py",
    ".ipynb",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".sql",
    ".sh",
    ".bash",
    ".zsh",
    ".xml",
    ".html",
    ".css",
    ".csv",
    ".tsv",
}

KEY_FILENAME_PRIORITY = {
    "readme.md",
    "readme.txt",
    "pyproject.toml",
    "requirements.txt",
    "setup.py",
    "package.json",
    "cargo.toml",
    "makefile",
    "dockerfile",
    "config.yaml",
    "config.yml",
    "main.py",
}

SOURCE_BLOG_TEXT_SUFFIXES = {
    ".md",
    ".txt",
}

SOURCE_BLOG_DOCUMENT_SUFFIXES = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
}

SOURCE_BLOG_STRUCTURED_SUFFIXES = {
    ".csv",
    ".tsv",
    ".json",
    ".jsonl",
}

SOURCE_BLOG_DATABASE_SUFFIXES = {
    ".db",
    ".sqlite",
    ".sqlite3",
}

SOURCE_BLOG_CONTROL_FILES = {
    "sources_url.json",
    "source_annotations.json",
}

def _source_metadata_field_specs() -> List[tuple[str, Any]]:
    return [
        ("Source Input", "source_input"),
        ("Source Type", "source_type"),
        ("Resolved From", lambda entry: entry.get("resolved_from") or ""),
        ("Local Path", "local_path"),
        ("User Note", lambda entry: entry.get("user_note") or ""),
        ("Prepared Markdown", "prepared_md_path"),
        ("Parse Status", lambda entry: entry.get("parse_status", "unknown")),
        ("File Count", lambda entry: entry.get("snapshot", {}).get("file_count", 0)),
    ]

def _generate_markdown_with_agent(prompt: str, llm_config: Dict[str, Any], agent_name: str) -> str:
    agent = ConversableAgent(
        name=agent_name,
        human_input_mode="NEVER",
        system_message="",
        llm_config=llm_config,
    )
    reply = agent.generate_reply(messages=[{"role": "user", "content": prompt}], sender=None)
    if isinstance(reply, dict):
        content = str(reply.get("content", "")).strip()
    else:
        content = str(reply).strip()
    content = re.sub(r"^```[a-zA-Z0-9_-]*\s*\n", "", content)
    content = re.sub(r"\n?```$", "", content).strip()
    return content


def _build_source_hint(entry: Dict[str, Any]) -> Optional[str]:
    source_type = entry.get("source_type")
    source_input = entry.get("source_input")
    if not source_input:
        return None
    return f"{source_type or 'source'}: {source_input}"


def _build_source_metadata_appendix(items: List[Dict[str, Any]]) -> str:
    lines = ["## Source Metadata Appendix", ""]
    for idx, entry in enumerate(items, start=1):
        lines.extend(
            [
                f"### Source {idx}",
                f"- item_id: `{entry.get('item_id', '')}`",
                f"- source_input: `{entry.get('source_input', '')}`",
                f"- source_type: `{entry.get('source_type', '')}`",
                f"- local_path: `{entry.get('local_path', '')}`",
                f"- user_note: `{entry.get('user_note', '')}`",
                f"- prepared_md_path: `{entry.get('prepared_md_path', '')}`",
                f"- blog_path: `{entry.get('blog_path', '')}`",
                f"- parse_status: `{entry.get('parse_status', 'unknown')}`",
                f"- blog_status: `{entry.get('blog_status', 'unknown')}`",
                "",
            ]
        )
    return "\n".join(lines).strip()


def _load_source_annotations(input_dir: Path, workspace_dir: Path) -> Dict[str, Dict[str, str]]:
    annotations_path = input_dir / "source_annotations.json"
    if not annotations_path.exists() or not annotations_path.is_file():
        return {}

    payload = json.loads(annotations_path.read_text(encoding="utf-8"))
    items = payload.get("items", [])
    if not isinstance(items, list):
        raise WorkflowError("source_annotations.json must contain a list under 'items'")

    annotations: Dict[str, Dict[str, str]] = {}
    for item in items:
        if not isinstance(item, dict):
            continue
        source = str(item.get("source", "")).strip()
        if not source:
            continue

        note = str(item.get("note", "")).strip()
        entry = {"note": note}
        for key in _build_source_annotation_keys(source, workspace_dir=workspace_dir, input_dir=input_dir):
            annotations[key] = entry

    return annotations


def _build_source_annotation_keys(source: str, workspace_dir: Path, input_dir: Path) -> List[str]:
    normalized = str(source).strip()
    if not normalized:
        return []

    keys: List[str] = [normalized]
    parsed = urlparse(normalized)
    if parsed.scheme in {"http", "https", "git", "ssh"} or normalized.startswith("git@"):
        return keys

    try:
        local_path = _resolve_local_path(normalized, workspace_dir)
        keys.append(str(local_path))
        try:
            keys.append(get_relative_path(local_path, workspace_dir))
        except Exception:
            pass
        try:
            keys.append(get_relative_path(local_path, input_dir))
        except Exception:
            pass
    except Exception:
        path_obj = Path(normalized)
        keys.append(path_obj.as_posix())
        if not path_obj.is_absolute():
            keys.append((workspace_dir / path_obj).as_posix())

    # Keep insertion order while removing duplicates.
    return list(dict.fromkeys(keys))


def _find_source_annotation(
    source: str,
    annotations: Dict[str, Dict[str, str]],
    workspace_dir: Path,
    input_dir: Path,
) -> Optional[Dict[str, str]]:
    if not annotations:
        return None

    for key in _build_source_annotation_keys(source, workspace_dir=workspace_dir, input_dir=input_dir):
        if key in annotations:
            return annotations[key]
    return None


def _resolve_source_input(
    source: str,
    workspace_dir: Path,
    bundle_dir: Path,
    timeout_seconds: Optional[int],
    max_total_bytes: Optional[int],
) -> Dict[str, Any]:
    if _looks_like_existing_local_path(source, workspace_dir):
        local_path = _resolve_local_path(source, workspace_dir)
        return {
            "source_type": "local_path",
            "local_path": str(local_path),
        }

    if _looks_like_git_url(source):
        local_path = _clone_git_repo(source, bundle_dir / "repo", timeout_seconds=timeout_seconds)
        return {
            "source_type": "git_url",
            "resolved_from": source,
            "local_path": str(local_path),
        }

    if _looks_like_http_url(source):
        local_path = _download_url(
            source,
            destination_dir=bundle_dir / "downloads",
            timeout_seconds=timeout_seconds,
            max_total_bytes=max_total_bytes,
        )
        return {
            "source_type": "url",
            "resolved_from": source,
            "local_path": str(local_path),
        }

    raise FileNotFoundError(f"Unsupported or unresolved source input: {source}")


def _looks_like_existing_local_path(source: str, workspace_dir: Path) -> bool:
    try:
        _resolve_local_path(source, workspace_dir)
        return True
    except FileNotFoundError:
        return False


def _looks_like_http_url(source: str) -> bool:
    parsed = urlparse(source)
    return parsed.scheme in {"http", "https"}


def _looks_like_git_url(source: str) -> bool:
    if source.startswith("git@"):
        return True
    parsed = urlparse(source)
    if parsed.scheme in {"git", "ssh"}:
        return True
    return source.endswith(".git")


def _prepare_source_bundle(
    item_meta: Dict[str, Any],
    workspace_dir: Path,
    max_files: Optional[int],
    max_depth: Optional[int],
    max_bytes_per_file: Optional[int],
    max_total_bytes_per_item: Optional[int],
    max_preview_rows: Optional[int],
) -> Dict[str, Any]:
    root_path = Path(str(item_meta["local_path"])).resolve()
    bundle_dir = workspace_dir / str(item_meta["bundle_dir"])
    bundle_dir.mkdir(parents=True, exist_ok=True)
    prepared_md_path = bundle_dir / "source.md"

    snapshot = _collect_path_snapshot(root_path, max_files=max_files, max_depth=max_depth)
    key_files = _pick_key_files(snapshot.get("files", []), limit=8)

    if root_path.is_file():
        prepared = _prepare_single_source_file_markdown(
            file_path=root_path,
            prepared_dir=prepared_md_path.parent,
            max_bytes=max_bytes_per_file,
            max_rows=max_preview_rows,
        )
        parse_status = "prepared_file"
    else:
        prepared = _prepare_source_tree_markdown(
            root_path=root_path,
            prepared_dir=prepared_md_path.parent,
            key_files=key_files,
            max_bytes_per_file=max_bytes_per_file,
            max_total_bytes=max_total_bytes_per_item,
            max_rows=max_preview_rows,
        )
        parse_status = "prepared_directory"

    header_lines = [
        "# Prepared Source",
        "",
        f"- Source Input: `{item_meta.get('source_input', '')}`",
        f"- Source Type: `{item_meta.get('source_type', '')}`",
        f"- Local Path: `{root_path}`",
        f"- Snapshot Kind: `{snapshot.get('kind', 'unknown')}`",
        f"- File Count: `{snapshot.get('file_count', 0)}`",
        f"- Snapshot Truncated: `{snapshot.get('truncated', False)}`",
        "",
        "## Tree Snapshot",
        "",
        "```text",
        "\n".join(snapshot.get("tree", [])) or "(empty)",
        "```",
        "",
        prepared["markdown"].strip(),
        "",
    ]
    save_markdown("\n".join(header_lines), prepared_md_path)

    return {
        "parse_status": parse_status,
        "prepared_md_path": get_relative_path(prepared_md_path, workspace_dir),
        "snapshot": {
            "kind": snapshot.get("kind"),
            "file_count": snapshot.get("file_count", 0),
            "truncated": snapshot.get("truncated", False),
            "tree": snapshot.get("tree", []),
            "key_files": key_files,
        },
        "preview_kind": prepared.get("preview_kind"),
        "bytes_used": prepared.get("bytes_used", 0),
    }


def _build_source_blog_prompt(entry: Dict[str, Any], workspace_dir: Path) -> str:
    prepared_md_path = workspace_dir / str(entry["prepared_md_path"])
    source_markdown = load_markdown(prepared_md_path) or ""
    available_images = (
        "If the prepared markdown already contains embedded markdown images, retain only the ones "
        "that are clearly important and preserve the original markdown image paths."
        if "![" in source_markdown
        else "No extracted images are available in this prepared markdown."
    )
    source_metadata = json.dumps(
        {
            "item_id": entry.get("item_id"),
            "source_input": entry.get("source_input"),
            "source_type": entry.get("source_type"),
            "resolved_from": entry.get("resolved_from"),
            "local_path": entry.get("local_path"),
            "user_note": entry.get("user_note"),
            "prepared_md_path": entry.get("prepared_md_path"),
            "snapshot": entry.get("snapshot", {}),
            "preview_kind": entry.get("preview_kind"),
            "parse_status": entry.get("parse_status"),
        },
        ensure_ascii=False,
        indent=2,
    )
    return SOURCE_BLOG_PROMPT.format(
        source_metadata=source_metadata,
        source_markdown=source_markdown,
        available_images=available_images,
    )


def _prepare_single_source_file_markdown(
    file_path: Path,
    prepared_dir: Path,
    max_bytes: Optional[int],
    max_rows: Optional[int],
) -> Dict[str, Any]:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _prepare_pdf_markdown_fragment(
            file_path=file_path,
            prepared_dir=prepared_dir,
            section_title="## Source PDF",
            max_bytes=max_bytes,
            use_shared_images_dir=True,
        )

    preview, preview_kind = _read_supported_preview(file_path, max_bytes=max_bytes, max_rows=max_rows)
    if preview is None or preview_kind is None:
        raise ValueError(f"Unsupported source file type: {suffix or '(none)'} ({file_path.name})")

    markdown = _format_source_preview_markdown(
        file_path=file_path,
        preview=preview,
        preview_kind=preview_kind,
        target_dir=prepared_dir,
    )
    return {
        "markdown": markdown,
        "preview_kind": preview_kind,
        "bytes_used": _estimate_preview_bytes(file_path, max_bytes=max_bytes),
    }


def _prepare_source_tree_markdown(
    root_path: Path,
    prepared_dir: Path,
    key_files: List[str],
    max_bytes_per_file: Optional[int],
    max_total_bytes: Optional[int],
    max_rows: Optional[int],
) -> Dict[str, Any]:
    used_bytes = 0
    sections = [
        "## High-Signal Files",
        "",
    ]

    if not key_files:
        sections.append("_No high-signal files were selected from this directory._")

    for relative_path in key_files:
        file_path = root_path / relative_path
        remaining = None if max_total_bytes is None else max_total_bytes - used_bytes
        if remaining is not None and remaining <= 0:
            sections.extend(
                [
                    f"### `{relative_path}`",
                    "",
                    "_Skipped because the per-item read budget was exhausted._",
                    "",
                ]
            )
            continue

        effective_limit = max_bytes_per_file
        if remaining is not None:
            effective_limit = remaining if effective_limit is None else min(effective_limit, remaining)

        if file_path.suffix.lower() == ".pdf":
            pdf_fragment = _prepare_pdf_markdown_fragment(
                file_path=file_path,
                prepared_dir=prepared_dir,
                section_title=f"### `{relative_path}`",
                max_bytes=effective_limit,
            )
            used_bytes += pdf_fragment["bytes_used"]
            sections.extend(
                [
                    pdf_fragment["markdown"],
                    "",
                ]
            )
            continue

        preview, preview_kind = _read_supported_preview(file_path, max_bytes=effective_limit, max_rows=max_rows)
        if preview is None or preview_kind is None:
            sections.extend(
                [
                    f"### `{relative_path}`",
                    "",
                    f"_Skipped unsupported file type: `{file_path.suffix.lower() or '(none)'}`._",
                    "",
                ]
            )
            continue

        used_bytes += _estimate_preview_bytes(file_path, max_bytes=effective_limit)
        sections.extend(
            [
                f"### `{relative_path}`",
                "",
                _format_source_preview_markdown(
                    file_path=file_path,
                    preview=preview,
                    preview_kind=preview_kind,
                    target_dir=prepared_dir,
                ),
                "",
            ]
        )

    return {
        "markdown": "\n".join(sections).strip(),
        "preview_kind": "directory",
        "bytes_used": used_bytes,
    }


def _prepare_pdf_markdown_fragment(
    file_path: Path,
    prepared_dir: Path,
    section_title: str,
    max_bytes: Optional[int],
    use_shared_images_dir: bool = False,
) -> Dict[str, Any]:
    if use_shared_images_dir:
        images_dir = prepared_dir / "images"
    else:
        assets_dir = prepared_dir / "_pdf_assets" / _safe_asset_id(file_path)
        images_dir = assets_dir / "images"
    pdf_parse = extract_pdf_markdown_with_images(
        file_path,
        images_dir=images_dir,
        markdown_dir=prepared_dir,
    )

    if pdf_parse["parse_status"] == "fulltext":
        markdown = "\n".join(
            [
                section_title,
                "",
                f"- Preview Kind: `document`",
                f"- Source PDF: `{file_path.name}`",
                f"- Parser Used: `{pdf_parse['parser_used']}`",
                f"- Extracted Images: `{pdf_parse['image_count']}`",
                "",
                pdf_parse["markdown_body"],
            ]
        )
    else:
        raise RuntimeError(f"PDF parsing failed for {file_path}")

    return {
        "markdown": markdown,
        "preview_kind": "document",
        "bytes_used": _estimate_preview_bytes(file_path, max_bytes=max_bytes),
    }


def _read_supported_preview(
    file_path: Path,
    max_bytes: Optional[int],
    max_rows: Optional[int],
) -> tuple[Optional[Dict[str, Any]], Optional[str]]:
    suffix = file_path.suffix.lower()

    if suffix in SOURCE_BLOG_TEXT_SUFFIXES:
        preview = _read_text_preview(file_path, max_bytes=max_bytes)
        return preview, "markdown" if suffix == ".md" else "text"
    if suffix in SOURCE_BLOG_DOCUMENT_SUFFIXES:
        if suffix == ".docx":
            return _read_docx_text(file_path, max_bytes=max_bytes), "document"
        if suffix == ".pptx":
            return _read_pptx_text(file_path, max_bytes=max_bytes), "document"
        return _read_excel_text(file_path, max_bytes=max_bytes), "document"
    if suffix in SOURCE_BLOG_STRUCTURED_SUFFIXES:
        return _preview_structured_file(file_path, max_rows=max_rows, max_bytes=max_bytes), "structured"
    if suffix in SOURCE_BLOG_DATABASE_SUFFIXES:
        return _preview_database_file(file_path, max_rows=max_rows), "database"
    return None, None


def _format_source_preview_markdown(
    file_path: Path,
    preview: Dict[str, Any],
    preview_kind: str,
    target_dir: Path,
) -> str:
    if preview_kind == "markdown":
        content = (preview.get("content") or "").strip()
        if content:
            content = rewrite_markdown_paths(
                content,
                source_dir=file_path.parent,
                target_dir=target_dir,
            )
        return content or "_No markdown content extracted from this source file._"

    if preview_kind == "structured":
        preview_json = json.dumps(preview, ensure_ascii=False, indent=2, default=str)
        return "\n".join(
            [
                "- Preview Kind: `structured`",
                "- Structured payload preview:",
                "",
                "```json",
                preview_json,
                "```",
            ]
        )

    if preview_kind == "database":
        preview_json = json.dumps(preview, ensure_ascii=False, indent=2, default=str)
        return "\n".join(
            [
                "- Preview Kind: `database`",
                "- Database payload preview:",
                "",
                "```json",
                preview_json,
                "```",
            ]
        )

    content = (preview.get("content") or "").strip()
    return "\n".join(
        [
            f"- Preview Kind: `{preview_kind}`",
            f"- Truncated: `{preview.get('truncated', False)}`",
            "",
            content or "_No text extracted from this source file._",
        ]
    )


def _estimate_preview_bytes(file_path: Path, max_bytes: Optional[int]) -> int:
    size = file_path.stat().st_size
    return size if max_bytes is None else min(size, max_bytes)


def _safe_asset_id(file_path: Path) -> str:
    digest = hashlib.sha1(str(file_path).encode("utf-8")).hexdigest()[:8]
    stem = re.sub(r"[^a-zA-Z0-9._-]+", "_", file_path.stem).strip("._")[:32] or "asset"
    return f"{stem}_{digest}"


def _safe_item_id(index: int, source: str) -> str:
    # Use a deterministic hash of the source input to generate a unique ID for each item.
    parsed = urlparse(source)
    if parsed.scheme and (parsed.path or parsed.netloc):
        name_hint = Path(parsed.path).name or parsed.netloc
    else:
        name_hint = Path(source).name or source

    if name_hint.endswith(".git"):
        name_hint = name_hint[:-4]

    sanitized_name = re.sub(r"[^a-zA-Z0-9_-]+", "_", name_hint).strip("._-").lower()
    if not sanitized_name:
        sanitized_name = "source"
    sanitized_name = sanitized_name[:32]

    digest = hashlib.sha1(source.encode("utf-8")).hexdigest()[:12]
    return f"{index:03d}_{sanitized_name}_{digest}"


def _resolve_local_path(source: str, workspace_dir: Path) -> Path:
    source_path = Path(source).expanduser()

    if source_path.is_absolute() and source_path.exists():
        return source_path.resolve()

    workspace_candidate = (workspace_dir / source).expanduser()
    if workspace_candidate.exists():
        return workspace_candidate.resolve()

    cwd_candidate = (Path.cwd() / source).expanduser()
    if cwd_candidate.exists():
        return cwd_candidate.resolve()

    if source_path.exists():
        return source_path.resolve()

    raise FileNotFoundError(f"Local path not found: {source}")


def _clone_git_repo(url: str, destination: Path, timeout_seconds: Optional[int]) -> Path:
    destination.parent.mkdir(parents=True, exist_ok=True)

    cmd = ["git", "clone", "--depth", "1", url, str(destination)]
    result = subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        timeout=timeout_seconds,
        check=False,
    )
    if result.returncode != 0:
        stderr = (result.stderr or "").strip()
        raise RuntimeError(f"git clone failed: {stderr or 'unknown error'}")
    return destination.resolve()


def _download_url(url: str, destination_dir: Path, timeout_seconds: Optional[int], max_total_bytes: Optional[int]) -> Path:
    destination_dir.mkdir(parents=True, exist_ok=True)

    parsed = urlparse(url)
    filename = Path(parsed.path).name or "downloaded_file"
    name_suffix = Path(filename).suffix.lower()

    downloaded = 0
    with requests.get(url, stream=True, timeout=timeout_seconds) as response:
        response.raise_for_status()
        content_type = (response.headers.get("Content-Type") or "").lower()
        iterator = response.iter_content(chunk_size=8192)
        first_chunk = next(iterator, b"")

        if name_suffix:
            target_path = destination_dir / filename
        else:
            if "application/pdf" in content_type or first_chunk.startswith(b"%PDF-"):
                ext = ".pdf"
            elif "text/html" in content_type:
                ext = ".html"
            elif "application/json" in content_type:
                ext = ".json"
            elif "text/plain" in content_type:
                ext = ".txt"
            elif "application/zip" in content_type:
                ext = ".zip"
            else:
                ext = ".bin"
            target_path = destination_dir / f"{filename}{ext}"

        with open(target_path, "wb") as fp:
            if first_chunk:
                downloaded += len(first_chunk)
                if max_total_bytes is not None and downloaded > max_total_bytes:
                    target_path.unlink(missing_ok=True)
                    raise RuntimeError(f"download exceeded max_total_bytes={max_total_bytes}")
                fp.write(first_chunk)
            for chunk in iterator:
                if not chunk:
                    continue
                downloaded += len(chunk)
                if max_total_bytes is not None and downloaded > max_total_bytes:
                    target_path.unlink(missing_ok=True)
                    raise RuntimeError(f"download exceeded max_total_bytes={max_total_bytes}")
                fp.write(chunk)

    return target_path.resolve()


def _collect_path_snapshot(path: Path, max_files: Optional[int], max_depth: Optional[int]) -> Dict[str, Any]:
    if path.is_file():
        stat = path.stat()
        return {
            "kind": "file",
            "root": str(path),
            "file_count": 1,
            "truncated": False,
            "files": [
                {
                    "path": path.name,
                    "size": stat.st_size,
                    "suffix": path.suffix.lower(),
                }
            ],
            "tree": [path.name],
        }

    file_records: List[Dict[str, Any]] = []
    tree_lines: List[str] = []
    truncated = False

    for current_root, dirnames, filenames in os.walk(path):
        rel_root = Path(current_root).resolve().relative_to(path.resolve())
        depth = len(rel_root.parts)

        if max_depth is not None and depth > max_depth:
            dirnames[:] = []
            continue

        dirnames.sort()
        filenames.sort()

        display_root = "." if str(rel_root) == "." else str(rel_root)
        tree_lines.append(f"{display_root}/")

        if max_depth is not None and depth >= max_depth:
            dirnames[:] = []

        for filename in filenames:
            absolute_file = Path(current_root) / filename
            try:
                stat = absolute_file.stat()
                size = stat.st_size
            except OSError:
                size = None

            rel_path = absolute_file.resolve().relative_to(path.resolve())
            file_records.append(
                {
                    "path": str(rel_path),
                    "size": size,
                    "suffix": absolute_file.suffix.lower(),
                    "name": absolute_file.name,
                }
            )
            tree_lines.append(f"  - {rel_path}")

            if max_files is not None and len(file_records) >= max_files:
                truncated = True
                break

        if truncated:
            break

    return {
        "kind": "directory",
        "root": str(path),
        "file_count": len(file_records),
        "truncated": truncated,
        "files": file_records,
        "tree": tree_lines,
    }


def _pick_key_files(file_records: List[Dict[str, Any]], limit: int = 8) -> List[str]:
    if not file_records:
        return []

    def score(record: Dict[str, Any]) -> tuple:
        name = (record.get("name") or "").lower()
        suffix = (record.get("suffix") or "").lower()
        is_priority_name = 1 if name in KEY_FILENAME_PRIORITY else 0
        is_text = 1 if suffix in TEXT_SUFFIXES else 0
        short_depth = -str(record.get("path", "")).count("/")
        return (is_priority_name, is_text, short_depth)

    ranked = sorted(file_records, key=score, reverse=True)
    return [r["path"] for r in ranked[:limit]]


def _read_text_preview(file_path: Path, max_bytes: Optional[int]) -> Dict[str, Any]:
    file_size = file_path.stat().st_size
    with open(file_path, "rb") as fp:
        payload = fp.read() if max_bytes is None else fp.read(max_bytes)
    truncated = False if max_bytes is None else file_size > max_bytes

    content = payload.decode("utf-8", errors="replace")

    return {
        "path": str(file_path),
        "size": file_size,
        "truncated": truncated,
        "content": content,
    }


def _truncate_by_bytes(content: str, max_bytes: Optional[int]) -> tuple[str, bool]:
    if max_bytes is None:
        return content, False
    raw = content.encode("utf-8")
    if len(raw) <= max_bytes:
        return content, False
    trimmed = raw[:max_bytes].decode("utf-8", errors="ignore")
    return trimmed, True


def _read_docx_text(file_path: Path, max_bytes: Optional[int]) -> Dict[str, Any]:
    from docx import Document

    file_size = file_path.stat().st_size
    doc = Document(str(file_path))
    paragraphs = [para.text for para in doc.paragraphs if para.text]
    content = "\n".join(paragraphs).strip()
    content, truncated = _truncate_by_bytes(content, max_bytes)
    return {
        "path": str(file_path),
        "size": file_size,
        "truncated": truncated,
        "content": content,
    }


def _read_pptx_text(file_path: Path, max_bytes: Optional[int]) -> Dict[str, Any]:
    from pptx import Presentation

    file_size = file_path.stat().st_size
    presentation = Presentation(str(file_path))
    lines: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                lines.append(text)
    content = "\n".join(lines).strip()
    content, truncated = _truncate_by_bytes(content, max_bytes)
    return {
        "path": str(file_path),
        "size": file_size,
        "truncated": truncated,
        "content": content,
    }


def _read_excel_text(file_path: Path, max_bytes: Optional[int]) -> Dict[str, Any]:
    import pandas as pd

    file_size = file_path.stat().st_size
    data_frames = pd.read_excel(file_path, sheet_name=None)
    chunks: List[str] = []
    for sheet_name, df in data_frames.items():
        chunks.append(f"# Sheet: {sheet_name}")
        chunks.append(df.to_string(index=False))
    content = "\n\n".join(chunks).strip()
    content, truncated = _truncate_by_bytes(content, max_bytes)
    return {
        "path": str(file_path),
        "size": file_size,
        "truncated": truncated,
        "content": content,
    }


def _preview_database_file(file_path: Path, max_rows: Optional[int]) -> Dict[str, Any]:
    import sqlite3

    rows_limit = max_rows if isinstance(max_rows, int) and max_rows > 0 else 3
    tables_limit = 10
    table_previews: List[Dict[str, Any]] = []
    connection = None
    try:
        connection = sqlite3.connect(f"file:{file_path}?mode=ro", uri=True)
        cursor = connection.execute(
            "SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%' ORDER BY name"
        )
        table_names = [str(row[0]) for row in cursor.fetchall() if row and row[0]]
        selected_tables = table_names[:tables_limit]
        truncated_tables = len(table_names) > tables_limit

        for table_name in selected_tables:
            safe_name = table_name.replace('"', '""')
            columns_cursor = connection.execute(f'PRAGMA table_info("{safe_name}")')
            columns = [
                {
                    "name": col[1],
                    "type": col[2],
                    "not_null": bool(col[3]),
                    "default": col[4],
                    "pk": bool(col[5]),
                }
                for col in columns_cursor.fetchall()
            ]

            sample_cursor = connection.execute(f'SELECT * FROM "{safe_name}" LIMIT {rows_limit}')
            column_names = [desc[0] for desc in (sample_cursor.description or [])]
            sample_rows = []
            for row in sample_cursor.fetchall():
                sample_rows.append(
                    {
                        column_names[idx]: (str(value)[:200] + "..." if isinstance(value, str) and len(value) > 200 else value)
                        for idx, value in enumerate(row)
                    }
                )

            table_previews.append(
                {
                    "table": table_name,
                    "columns": columns,
                    "sample_rows": sample_rows,
                    "preview_rows": len(sample_rows),
                }
            )

        return {
            "path": str(file_path),
            "format": "sqlite",
            "table_count": len(table_names),
            "tables_truncated": truncated_tables,
            "tables": table_previews,
        }
    except Exception as exc:
        return {
            "path": str(file_path),
            "format": "sqlite",
            "table_count": 0,
            "tables_truncated": False,
            "tables": [],
            "error": str(exc),
        }
    finally:
        if connection is not None:
            connection.close()


def _preview_structured_file(file_path: Path, max_rows: Optional[int], max_bytes: Optional[int]) -> Dict[str, Any]:
    suffix = file_path.suffix.lower()

    if suffix in {".csv", ".tsv"}:
        delimiter = "," if suffix == ".csv" else "\t"
        with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as fp:
            reader = csv.DictReader(fp, delimiter=delimiter)
            rows = []
            for idx, row in enumerate(reader):
                if max_rows is not None and idx >= max_rows:
                    break
                rows.append(row)
            return {
                "path": str(file_path),
                "format": suffix[1:],
                "columns": reader.fieldnames or [],
                "rows": rows,
                "preview_rows": len(rows),
            }

    if suffix in {".json", ".jsonl"}:
        file_size = file_path.stat().st_size
        with open(file_path, "rb") as fp:
            payload = fp.read() if max_bytes is None else fp.read(max_bytes)
        truncated = False if max_bytes is None else file_size > max_bytes

        if suffix == ".jsonl":
            text = payload.decode("utf-8", errors="replace")
            lines = [line for line in text.splitlines() if line.strip()]
            if max_rows:
                lines = lines[:max_rows]
            parsed_rows = []
            for line in lines:
                try:
                    parsed_rows.append(json.loads(line))
                except Exception:
                    parsed_rows.append({"raw": line})
            return {
                "path": str(file_path),
                "format": "jsonl",
                "rows": parsed_rows,
                "preview_rows": len(parsed_rows),
                "truncated": truncated,
            }

        text = payload.decode("utf-8", errors="replace")
        if truncated:
            return {
                "path": str(file_path),
                "format": "json",
                "truncated": True,
                "note": "File too large for full JSON parsing within max_bytes limit",
                "preview_text": text[: min(len(text), 2000)],
            }

        data = json.loads(text)
        if isinstance(data, dict):
            keys = list(data.keys())
            if max_rows:
                keys = keys[: max_rows * 2]
                sample = {k: data[k] for k in keys[:max_rows]}
            else:
                sample = data
            return {
                "path": str(file_path),
                "format": "json",
                "kind": "object",
                "keys": keys,
                "sample": sample,
                "truncated": truncated,
            }

        if isinstance(data, list):
            return {
                "path": str(file_path),
                "format": "json",
                "kind": "array",
                "length": len(data),
                "sample": data[:max_rows] if max_rows else data,
                "truncated": truncated,
            }

        return {
            "path": str(file_path),
            "format": "json",
            "kind": type(data).__name__,
            "value": data,
            "truncated": truncated,
        }

    raise ValueError(f"Unsupported structured file type: {suffix}")
